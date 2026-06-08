import io
import re
from html.parser import HTMLParser

_DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
_XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
_XLS_MIME = "application/vnd.ms-excel"
# Keep in sync with filetypes._CALENDAR (separate module, no shared import).
_CALENDAR_MIMES = {"text/calendar", "application/ics", "text/x-vcalendar"}
# Keep in sync with filetypes._CONTACTS (separate module, no shared import).
_VCARD_MIMES = {"text/x-vcard", "text/vcard", "application/vcard", "text/directory"}
_TEXTLIKE_APP_MIMES = {
    "application/json", "application/xml",
    "application/x-yaml", "application/yaml",
}


class _TextExtractor(HTMLParser):
    def __init__(self):
        super().__init__()
        self.parts = []
        self._skip = False

    def handle_starttag(self, tag, attrs):
        if tag in ("script", "style"):
            self._skip = True

    def handle_endtag(self, tag):
        if tag in ("script", "style"):
            self._skip = False

    def handle_data(self, data):
        if not self._skip:
            self.parts.append(data)


def html_to_text(html: str) -> str:
    p = _TextExtractor()
    p.feed(html or "")
    return " ".join("".join(p.parts).split())


def _pdf_text(data: bytes) -> str:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(data))
    return "\n".join((page.extract_text() or "") for page in reader.pages)


def _docx_text(data: bytes) -> str:
    import docx
    d = docx.Document(io.BytesIO(data))
    return "\n".join(par.text for par in d.paragraphs)


def _pptx_text(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
            elif shape.has_table:  # text-frame check misses table cells
                for trow in shape.table.rows:
                    parts.append("\t".join(cell.text for cell in trow.cells))
    return "\n".join(p for p in parts if p and p.strip())


def _cell_str(c) -> str:
    # Spreadsheet libs hand back whole numbers as floats (999 -> 999.0); render
    # those as plain integers so the text reads naturally and searches cleanly.
    if isinstance(c, float) and c.is_integer():
        return str(int(c))
    return str(c)


def _rows_text(rows) -> str:
    out = []
    for row in rows:
        cells = [s for s in (_cell_str(c) for c in row if c is not None) if s.strip()]
        if cells:
            out.append("\t".join(cells))
    return "\n".join(out)


def _xlsx_text(data: bytes) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    try:
        return "\n".join(_rows_text(ws.iter_rows(values_only=True)) for ws in wb.worksheets)
    finally:
        wb.close()


def _xls_text(data: bytes) -> str:
    import xlrd
    book = xlrd.open_workbook(file_contents=data)
    sheets = []
    for sheet in book.sheets():
        rows = (sheet.row_values(r) for r in range(sheet.nrows))
        sheets.append(_rows_text(rows))
    return "\n".join(sheets)


def _ics_unescape(v: str) -> str:
    return (v.replace("\\N", "\n").replace("\\n", "\n")
             .replace("\\,", ",").replace("\\;", ";").replace("\\\\", "\\"))


def _ics_when(v: str) -> str:
    m = re.match(r"^(\d{4})(\d{2})(\d{2})(?:T(\d{2})(\d{2})(\d{2})(Z)?)?$", v.strip())
    if not m:
        return v
    y, mo, d, h, mi, _s, z = m.groups()
    if h is None:
        return "%s-%s-%s" % (y, mo, d)
    return "%s-%s-%s %s:%s%s" % (y, mo, d, h, mi, " UTC" if z else "")


def _ics_text(data: bytes) -> str:
    raw = data.decode("utf-8", "replace")
    # RFC 5545 line unfolding: a line starting with space/tab continues the previous one.
    lines = []
    for line in raw.splitlines():
        if line[:1] in (" ", "\t") and lines:
            lines[-1] += line[1:]
        else:
            lines.append(line)
    events = []
    cur = None
    for line in lines:
        key_line = line.strip()
        if key_line == "BEGIN:VEVENT":
            cur = {"attendees": []}
        elif key_line == "END:VEVENT":
            if cur is not None:
                events.append(cur)
                cur = None
        elif cur is not None and ":" in line:
            name, value = line.split(":", 1)
            key = name.split(";", 1)[0].upper()
            if key == "SUMMARY":
                cur["summary"] = _ics_unescape(value)
            elif key == "DTSTART":
                cur["start"] = _ics_when(value)
            elif key == "DTEND":
                cur["end"] = _ics_when(value)
            elif key == "LOCATION":
                cur["location"] = _ics_unescape(value)
            elif key == "ORGANIZER":
                cur["organizer"] = value.replace("mailto:", "").replace("MAILTO:", "")
            elif key == "ATTENDEE":
                cur["attendees"].append(value.replace("mailto:", "").replace("MAILTO:", ""))
            elif key == "DESCRIPTION":
                cur["description"] = _ics_unescape(value)
    out = []
    for e in events:
        if e.get("summary"):
            out.append("Summary: " + e["summary"])
        when = e.get("start", "")
        if e.get("end"):
            when = (when + " → " + e["end"]) if when else e["end"]
        if when:
            out.append("When: " + when)
        if e.get("location"):
            out.append("Location: " + e["location"])
        if e.get("organizer"):
            out.append("Organizer: " + e["organizer"])
        if e.get("attendees"):
            out.append("Attendees: " + ", ".join(e["attendees"]))
        if e.get("description"):
            out.append("Description: " + e["description"])
        out.append("")
    return "\n".join(out).strip()


def _vcard_text(data: bytes) -> str:
    raw = data.decode("utf-8", "replace")
    lines = []
    for line in raw.splitlines():
        if line[:1] in (" ", "\t") and lines:
            lines[-1] += line[1:]
        else:
            lines.append(line)
    cards = []
    cur = None
    for line in lines:
        u = line.strip().upper()
        if u == "BEGIN:VCARD":
            cur = {"emails": [], "tels": []}
        elif u == "END:VCARD":
            if cur is not None:
                cards.append(cur)
                cur = None
        elif cur is not None and ":" in line:
            name, value = line.split(":", 1)
            key = name.split(";", 1)[0].upper()
            value = _ics_unescape(value)
            if key == "FN":
                cur["name"] = value
            elif key == "N" and "name" not in cur:
                cur["name"] = value.replace(";", " ").strip()
            elif key == "EMAIL":
                cur["emails"].append(value)
            elif key == "TEL":
                cur["tels"].append(value)
            elif key == "ORG":
                cur["org"] = value.replace(";", " ").strip()
            elif key == "TITLE":
                cur["title"] = value
            elif key == "ADR":
                cur["adr"] = value.replace(";", " ").strip()
            elif key == "URL":
                cur["url"] = value
    out = []
    for c in cards:
        if c.get("name"):
            out.append("Name: " + c["name"])
        if c.get("title"):
            out.append("Title: " + c["title"])
        if c.get("org"):
            out.append("Organization: " + c["org"])
        if c.get("emails"):
            out.append("Email: " + ", ".join(c["emails"]))
        if c.get("tels"):
            out.append("Phone: " + ", ".join(c["tels"]))
        if c.get("adr"):
            out.append("Address: " + c["adr"])
        if c.get("url"):
            out.append("URL: " + c["url"])
        out.append("")
    return "\n".join(out).strip()


def extract_text(filename: str, mime: str, data: bytes) -> str:
    """Best-effort plain-text extraction. Returns '' for unsupported or on any error."""
    mime = (mime or "").lower()
    try:
        if mime == "application/pdf":
            return _pdf_text(data)
        if mime == _DOCX_MIME:
            return _docx_text(data)
        if mime == _PPTX_MIME:
            return _pptx_text(data)
        if mime == _XLSX_MIME:
            return _xlsx_text(data)
        if mime == _XLS_MIME:
            return _xls_text(data)
        if mime in _CALENDAR_MIMES:
            return _ics_text(data)
        if mime in _VCARD_MIMES:
            return _vcard_text(data)
        if mime in _TEXTLIKE_APP_MIMES:
            return data.decode("utf-8", "replace")
        if mime.startswith("text/html"):
            return html_to_text(data.decode("utf-8", "replace"))
        if mime.startswith("text/"):
            return data.decode("utf-8", "replace")
    except Exception:
        return ""
    return ""
