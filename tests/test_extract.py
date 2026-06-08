import io
import docx
from reportlab.pdfgen import canvas
from mboxviewer.extract import extract_text, html_to_text


def _pdf(text):
    buf = io.BytesIO()
    c = canvas.Canvas(buf)
    c.drawString(72, 720, text)
    c.save()
    return buf.getvalue()


def _docx(text):
    buf = io.BytesIO()
    d = docx.Document()
    d.add_paragraph(text)
    d.save(buf)
    return buf.getvalue()


def test_extract_plain_text():
    assert "hello" in extract_text("a.txt", "text/plain", b"hello world")


def test_extract_pdf():
    out = extract_text("a.pdf", "application/pdf", _pdf("INVOICE 12345"))
    assert "12345" in out


def test_extract_docx():
    mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    out = extract_text("a.docx", mime, _docx("QUARTERLY REPORT"))
    assert "QUARTERLY" in out


def test_extract_unsupported_returns_empty():
    assert extract_text("a.bin", "application/octet-stream", b"\x00\x01") == ""


def test_extract_handles_corrupt_pdf_gracefully():
    assert extract_text("a.pdf", "application/pdf", b"not a real pdf") == ""


def test_html_to_text_strips_tags():
    assert html_to_text("<p>Hello <b>world</b></p>").strip() == "Hello world"


def test_extract_pptx():
    import io
    from pptx import Presentation
    from pptx.util import Inches
    from mboxviewer.extract import extract_text
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "ROADMAP Q3 LAUNCH"
    buf = io.BytesIO(); prs.save(buf)
    mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    assert "ROADMAP Q3 LAUNCH" in extract_text("d.pptx", mime, buf.getvalue())


def test_extract_xlsx():
    import io
    import openpyxl
    from mboxviewer.extract import extract_text
    wb = openpyxl.Workbook(); ws = wb.active
    ws["A1"] = "Region"; ws["B1"] = "Sales"; ws["A2"] = "EMEA"; ws["B2"] = 4200
    buf = io.BytesIO(); wb.save(buf)
    mime = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    text = extract_text("d.xlsx", mime, buf.getvalue())
    assert "EMEA" in text and "4200" in text


def test_extract_xls():
    import io
    import xlwt
    from mboxviewer.extract import extract_text
    wb = xlwt.Workbook(); ws = wb.add_sheet("S1")
    ws.write(0, 0, "Account"); ws.write(0, 1, "Balance"); ws.write(1, 0, "ACME"); ws.write(1, 1, 999)
    buf = io.BytesIO(); wb.save(buf)
    text = extract_text("d.xls", "application/vnd.ms-excel", buf.getvalue())
    # whole numbers render as plain integers, not "999.0"
    assert "ACME" in text and "999" in text and "999.0" not in text


def test_extract_pptx_table():
    import io
    from pptx import Presentation
    from pptx.util import Inches
    from mboxviewer.extract import extract_text
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    table = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2)).table
    table.cell(0, 0).text = "MILESTONE"; table.cell(0, 1).text = "Q4 SHIP"
    buf = io.BytesIO(); prs.save(buf)
    mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    text = extract_text("d.pptx", mime, buf.getvalue())
    assert "MILESTONE" in text and "Q4 SHIP" in text


def test_extract_corrupt_office_returns_empty():
    from mboxviewer.extract import extract_text
    mime = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    assert extract_text("d.xlsx", mime, b"not a real zip") == ""


def test_extract_ics_event():
    from mboxviewer.extract import extract_text
    ics = (
        "BEGIN:VCALENDAR\r\nBEGIN:VEVENT\r\n"
        "SUMMARY:Project Kickoff\r\n"
        "DTSTART:20240115T090000Z\r\nDTEND:20240115T100000Z\r\n"
        "LOCATION:Room 4\r\n"
        "ORGANIZER:mailto:alice@example.com\r\n"
        "ATTENDEE:mailto:bob@example.com\r\n"
        "DESCRIPTION:Discuss the plan\r\n"
        "END:VEVENT\r\nEND:VCALENDAR\r\n"
    )
    out = extract_text("invite.ics", "application/ics", ics.encode())
    assert "Project Kickoff" in out and "Room 4" in out
    assert "alice@example.com" in out and "bob@example.com" in out
    assert "2024-01-15" in out  # DTSTART is pretty-printed


def test_extract_ics_unfolds_long_lines():
    from mboxviewer.extract import extract_text
    ics = (
        "BEGIN:VEVENT\r\n"
        "SUMMARY:Quarterly planning and budget \r\n review session\r\n"
        "END:VEVENT\r\n"
    )
    out = extract_text("x.ics", "text/calendar", ics.encode())
    assert "Quarterly planning and budget review session" in out


def test_extract_ics_unfolds_midword():
    # RFC 5545 folds at an arbitrary octet boundary; unfolding must NOT insert a
    # space. A mid-word fold ("Kick\r\n off") must rejoin to "Kickoff".
    from mboxviewer.extract import extract_text
    ics = "BEGIN:VEVENT\r\nSUMMARY:Kick\r\n off\r\nEND:VEVENT\r\n"
    out = extract_text("x.ics", "text/calendar", ics.encode())
    assert "Kickoff" in out


def test_extract_textlike_application_types():
    from mboxviewer.extract import extract_text
    assert "hello" in extract_text("a.json", "application/json", b'{"k": "hello"}')
    assert "<note>" in extract_text("a.xml", "application/xml", b"<note>hi</note>")


def test_extract_ics_no_event_returns_empty():
    from mboxviewer.extract import extract_text
    assert extract_text("x.ics", "text/calendar", b"BEGIN:VCALENDAR\r\nEND:VCALENDAR\r\n") == ""


def test_extract_vcard():
    from mboxviewer.extract import extract_text
    vcf = ("BEGIN:VCARD\r\nVERSION:3.0\r\nFN:Alice Smith\r\nORG:Acme\r\n"
           "TITLE:Engineer\r\nEMAIL:alice@example.com\r\n"
           "EMAIL;TYPE=work:a.smith@acme.com\r\nTEL:+1-555-1234\r\nEND:VCARD\r\n")
    out = extract_text("c.vcf", "text/x-vcard", vcf.encode())
    assert "Name: Alice Smith" in out
    assert "alice@example.com" in out and "a.smith@acme.com" in out
    assert "+1-555-1234" in out and "Acme" in out


def test_extract_vcard_no_card_empty():
    from mboxviewer.extract import extract_text
    assert extract_text("x.vcf", "text/x-vcard", b"not a vcard") == ""


def test_doc_salvage_text_utf16():
    from mboxviewer.extract import _salvage_text
    raw = b"\x07\x00" + "Quarterly Report Growth".encode("utf-16-le") + b"\x00\x13"
    assert "Quarterly Report Growth" in _salvage_text(raw)


def test_doc_salvage_text_cp1252():
    from mboxviewer.extract import _salvage_text
    raw = b"\x00\x01" + "Hello cp1252 World".encode("cp1252") + b"\xff"
    assert "Hello cp1252 World" in _salvage_text(raw)


def test_extract_doc_non_ole_returns_empty():
    from mboxviewer.extract import extract_text
    assert extract_text("a.doc", "application/msword", b"not an ole file") == ""
